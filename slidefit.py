#! /usr/bin/env python3

# Copyright (C) 2020-2021 Tomas Nordin

# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.

# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
# GNU General Public License for more details.

# You should have received a copy of the GNU General Public License
# along with this program.  If not, see <https://www.gnu.org/licenses/>.

"""FIT IMAGES INTO SLIDES

Usage:
    slidefit [-d -m MATRIX -t PPT -o PPT]
             [ -y SPACE -x SPACE -p PAGES] IMG...
    slidefit (-h | -v)

Fit the image matrix onto each slide, keeping image aspect ratio. Assume
all images has the same ratio as the first image.

Images are added row-wise.

Options:
    -m, --matrix MATRIX
                  The number of rows and columns to put on each slide,
                  given as rows,columns. [default: 1,1]
    -o, --ppt PPT
                  Name of the output slide file (.pptx).
                  [default: slidefit-output.pptx]
    -t, --template PPT
                  Name of a possible input slide file.
    -p, --pages PAGES
                  The pages of a template file to use. For example
                  '2,5-7,9-12' to use the pages 2, 5 to 7 and 9 to 12.
                  The first slide is 1. This option is ignored if
                  --template is not given.
    -y, --vspace SPACE
                  Reserve a percentage of vertical space on the top
                  and/or bottom of the slide when adding images.
                  [default: 0,0]. (The sum of both should be < 100). Can
                  be floating numbers.
    -x, --hspace SPACE
                  Reserve a percentage of horizontal space on the left
                  and/or right of the slide when adding images.
                  [default: 0,0]. (The sum of both should be < 100). Can
                  be floating numbers.
    -h, --help    Show this help and exit
    -v, --version
                  Show version.
    -d, --debug   Show args and exit

Arguments:
    IMG           The images to merge into the slides. Merged in given
                  order with one MATRIX on each slide.

"""
import os
import math
from collections import namedtuple
import sys

import pptx
import docopt

__version__ = '0.1.1'

# 1 mm is 36000 English Metric Units (EMU)
# blank_slide_layout = prs.slide_layouts[6]:
# width [mm]: 254.0
# height [mm] 190.5

Matrix = namedtuple('Matrix', ('rows', 'cols'))
Size = namedtuple('Size', ('width', 'height'))
Location = namedtuple('Location', ('left', 'top'))
Space = namedtuple('Space', ('top', 'bottom', 'left', 'right'))


def split2ints(rangestring):
    # help function for the --pages option
    # decrement with 1, yield numbers inclusive
    for number in rangestring.split(','):
        try:
            yield int(number) - 1
        except ValueError:
            n1, n2 = [int(n) - 1 for n in number.split('-')]
            for n in range(n1, n2 + 1):
                yield n


def numberpair(args, option, floatp=True):
    # aid to provide useful error message
    numbers = args[option].split(',')
    try:
        if floatp:
            return [float(numbers[0]), float(numbers[1])]
        else:
            return [int(numbers[0]), int(numbers[1])]
    except (IndexError, ValueError):
        sys.exit('slidefit: Bad argument for %s "%s"' % (option, args[option]))


def main(args):

    matrix = Matrix(*numberpair(args, '--matrix', floatp=False))

    # dummy presentation to inspect resulting size
    prs = pptx.Presentation(args['--template'])
    if args['--template']:
        try:
            layout = prs.slides[-1].slide_layout
        except IndexError:
            layout = prs.slide_layouts[6]
    else:
        layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    pic = slide.shapes.add_picture(args['IMG'][0], 36000, 36000)
    orig_size = Size(pic.width, pic.height)  # emu

    reserved = Space(*([int(float(perc) / 100 * prs.slide_height)
                        for perc in numberpair(args, '--vspace',
                                               floatp=True)]
                       + [int(float(perc) / 100 * prs.slide_width)
                          for perc in numberpair(args, '--hspace',
                                                 floatp=True)]))

    height_avail = prs.slide_height - reserved.top - reserved.bottom
    width_avail = prs.slide_width - reserved.left - reserved.right

    # pic size:
    if orig_size.width * matrix.cols < width_avail:
        width = orig_size.width
    else:
        width = width_avail // matrix.cols

    if orig_size.height * matrix.rows < height_avail:
        height = orig_size.height
    else:
        height = height_avail // matrix.rows

    #  which one is the smaller fraction of origin?
    widthfrac = width / orig_size.width
    heightfrac = height / orig_size.height
    if widthfrac < heightfrac:  # width is limiting
        pic_size = Size(width, int(orig_size.height * widthfrac))
    else:
        pic_size = Size(int(orig_size.width * heightfrac), height)

    # calculate and make a sequence of locations for one slide matrix:
    widthmargin = width_avail - pic_size.width * matrix.cols
    heightmargin = height_avail - pic_size.height * matrix.rows
    widthfill = widthmargin // (matrix.cols + 1)
    heightfill = heightmargin // (matrix.rows + 1)
    lefts = [n * pic_size.width + widthfill + reserved.left
             for n in range(matrix.cols)]
    tops = [n * pic_size.height + heightfill + reserved.top
            for n in range(matrix.rows)]

    locations = []
    for top in tops:
        for left in lefts:
            locations.append(Location(left, top))

    reqslidescnt = math.ceil(len(args['IMG']) / (matrix.rows * matrix.cols))

    prs = pptx.Presentation(args['--template'])

    # add or select slides according to the command:
    if args['--template'] and not args['--pages']:
        # use slides already in template and add slides as required
        # layout is given by template last slide
        layout = prs.slides[-1].slide_layout
        slides = [slide for slide in prs.slides]
        missingcnt = reqslidescnt - len(slides)
        slides += [prs.slides.add_slide(layout) for n in range(missingcnt)]
    elif args['--template'] and args['--pages']:
        # slides are the slides that exist and correspond to --pages
        # not all images might be added.
        slideindexes = sorted(list(split2ints(args['--pages'])))
        slides = []
        for index in slideindexes:
            try:
                slides.append(prs.slides[index])
            except IndexError:
                print('Warning: pages missing in template relative --pages')
                break
    elif not args['--template']:
        layout = prs.slide_layouts[6]  # blank slide layout
        slides = [prs.slides.add_slide(layout) for n in range(reqslidescnt)]

    # add the images to selected slides
    for slide in slides:
        for location in locations:
            if not args['IMG']:
                break
            slide.shapes.add_picture(args['IMG'].pop(0),
                                     location.left, location.top,
                                     width=pic_size.width,
                                     height=pic_size.height)

    if args['IMG']:
        print('Warning: not all images could be inserted by given command')

    if not os.path.exists(args['--ppt']):
        prs.save(args['--ppt'])
        print('Saved', args['--ppt'])
    else:
        sys.exit('slidefit: File exist, not over-writing (%s)' % args['--ppt'])


def debug(args):
    print(args)
    if args['--pages']:
        print(list(split2ints(args['--pages'])))
    exit(0)


def distmain():

    args = docopt.docopt(__doc__, version=__version__)
    if args['--debug']:
        debug(args)
    main(args)


if __name__ == '__main__':
    distmain()
